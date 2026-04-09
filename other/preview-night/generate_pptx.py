"""Generate Preview Night PowerPoint presentation for CS 7180."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG_DIR = Path(__file__).parent / "images"

# Northeastern red
NEU_RED = RGBColor(0xCC, 0x00, 0x00)
NEU_DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT_BLUE = RGBColor(0x1E, 0x90, 0xFF)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=18,
                          color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                          line_spacing=1.5, font_name="Calibri"):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle tuples for (text, bold, color, size) overrides
        if isinstance(line, tuple):
            p.text = line[0]
            p.font.bold = line[1] if len(line) > 1 else bold
            p.font.color.rgb = line[2] if len(line) > 2 else color
            p.font.size = Pt(line[3] if len(line) > 3 else font_size)
        else:
            p.text = line
            p.font.bold = bold
            p.font.color.rgb = color
            p.font.size = Pt(font_size)

        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))

    return txBox


def add_red_accent_bar(slide):
    """Add the thin red accent bar at top of slide."""
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), NEU_RED)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_background(slide, NEU_DARK)

# Red accent bar at top
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), NEU_RED)

# Course code
add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.6),
            "CS 7180", font_size=24, color=NEU_RED, bold=True)

# Title
add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "Vibe Coding", font_size=60, color=WHITE, bold=True)

# Subtitle
add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
            "AI-Assisted Software Engineering", font_size=32, color=RGBColor(0xCC, 0xCC, 0xCC))

# Divider line
add_shape(slide, Inches(1), Inches(4.0), Inches(3), Inches(0.04), NEU_RED)

# Details
add_multiline_textbox(slide, Inches(1), Inches(4.3), Inches(11), Inches(2),
                      [
                          "Fall 2026  |  Silicon Valley Campus  |  Once per week",
                          "",
                          "John Alexis Guerra Gomez",
                          "jguerra@northeastern.edu",
                      ],
                      font_size=20, color=RGBColor(0xAA, 0xAA, 0xAA), line_spacing=1.4)

# Background image (behind all content — add first, then send to back)
img_path = IMG_DIR / "image_1.png"
if img_path.exists():
    pic = slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    # Send image to back so text is on top
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)  # after background

# ============================================================
# SLIDE 2: The Hook
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_red_accent_bar(slide)

# Title
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            'Not Just "Vibe Coding"', font_size=44, color=NEU_DARK, bold=True)

# Quote box
quote_shape = add_shape(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.2),
                        RGBColor(0xFD, 0xF0, 0xF0))
add_shape(slide, Inches(0.8), Inches(1.6), Inches(0.08), Inches(1.2), NEU_RED)

add_textbox(slide, Inches(1.2), Inches(1.75), Inches(10.8), Inches(0.9),
            'Everyone talks about vibe coding — we teach you to do it professionally',
            font_size=26, color=NEU_DARK, bold=False,
            font_name="Calibri")

# Key points
add_multiline_textbox(slide, Inches(0.8), Inches(3.2), Inches(11), Inches(3.5),
                      [
                          ("AI speed  +  engineering quality  =  industry-ready", True, NEU_DARK, 28),
                          "",
                          "▸  Not just prompting — building production software with AI",
                          "▸  Test-Driven Development, CI/CD pipelines, code evaluations",
                          "▸  The course the industry wishes every new hire had taken",
                      ],
                      font_size=22, color=MID_GRAY, line_spacing=1.6)

# Split image at bottom
img_path = IMG_DIR / "image_2.png"
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(1.5), Inches(5.2), Inches(10), Inches(2.1))

# ============================================================
# SLIDE 3: Three AI Modalities
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_red_accent_bar(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "Three Ways to Code with AI", font_size=44, color=NEU_DARK, bold=True)

# Three modalities image as background illustration
img_path = IMG_DIR / "image_3.png"
if img_path.exists():
    pic = slide.shapes.add_picture(str(img_path), Inches(0), Inches(1.5), SLIDE_W, Inches(4.8))
    # Send behind text content
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

# Three columns
col_width = Inches(3.5)
col_colors = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_PURPLE]
col_titles = ["Conversational AI", "AI-Powered IDE", "Agentic Coding"]
col_weeks = ["Weeks 4–5", "Weeks 6–8", "Weeks 10–14"]
col_desc = [
    "Architecture\nPrototyping\nBrainstorming",
    "Daily coding workflow\nAutocomplete\nInline editing",
    "Autonomous agents\nAutomation\nMulti-file refactoring",
]

for i in range(3):
    x = Inches(0.8 + i * 4.0)
    # Card background
    card = add_shape(slide, x, Inches(1.8), col_width, Inches(3.8), LIGHT_GRAY)
    # Color bar at top of card
    add_shape(slide, x, Inches(1.8), col_width, Inches(0.08), col_colors[i])
    # Title
    add_textbox(slide, x + Inches(0.3), Inches(2.1), col_width - Inches(0.6), Inches(0.6),
                col_titles[i], font_size=22, color=NEU_DARK, bold=True)
    # Weeks
    add_textbox(slide, x + Inches(0.3), Inches(2.7), col_width - Inches(0.6), Inches(0.4),
                col_weeks[i], font_size=16, color=col_colors[i], bold=True)
    # Description
    add_textbox(slide, x + Inches(0.3), Inches(3.2), col_width - Inches(0.6), Inches(2),
                col_desc[i], font_size=18, color=MID_GRAY)

# Bottom note
add_multiline_textbox(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(1.2),
                      [
                          ("Plus: Prompt engineering · Context engineering · Evals · Agent architectures", False, MID_GRAY, 18),
                          ("📌 Curriculum adapts to the latest AI tools and practices", True, NEU_RED, 18),
                      ],
                      font_size=18, color=MID_GRAY, line_spacing=1.6)

# ============================================================
# SLIDE 4: Course Topics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_red_accent_bar(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "15 Weeks of AI + Engineering", font_size=44, color=NEU_DARK, bold=True)

# Timeline entries
phases = [
    ("Weeks 1–3", "LLM fundamentals, prompt engineering, model comparison", ACCENT_BLUE),
    ("Weeks 4–5", "Conversational AI coding, rapid prototyping", ACCENT_BLUE),
    ("Weeks 6–8", "IDE AI tools, context engineering, MCP servers", ACCENT_GREEN),
    ("Weeks 10–12", "Agentic coding — skills, hooks, sub-agents, TDD", ACCENT_PURPLE),
    ("Weeks 13–14", "Agent architectures, Agent SDK, AI engineering", ACCENT_PURPLE),
]

for i, (week, desc, color) in enumerate(phases):
    y = Inches(1.7 + i * 0.85)
    # Color dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.1), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    # Week label
    add_textbox(slide, Inches(1.5), y, Inches(2.2), Inches(0.5),
                week, font_size=20, color=NEU_DARK, bold=True)
    # Description
    add_textbox(slide, Inches(3.8), y, Inches(8), Inches(0.5),
                desc, font_size=20, color=MID_GRAY)

# Vertical line connecting dots
add_shape(slide, Inches(1.1), Inches(2.05), Inches(0.04), Inches(3.2), RGBColor(0xDD, 0xDD, 0xDD))

# Tech stack banner
tech_bg = add_shape(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.7), LIGHT_GRAY)
add_textbox(slide, Inches(1.2), Inches(6.05), Inches(11), Inches(0.6),
            "Full-stack:  React/Next.js  ·  Node.js  ·  PostgreSQL/MongoDB  ·  GitHub Actions",
            font_size=20, color=NEU_DARK, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: Three Portfolio Projects
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_red_accent_bar(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "3 Portfolio-Ready Applications", font_size=44, color=NEU_DARK, bold=True)

projects = [
    ("Project 1", "Personal Utility App", "solo",
     "Solve a real, user-validated problem\nwith conversational AI", NEU_RED),
    ("Project 2", "Full-Stack Application", "pair",
     "Auth, public API, 80%+ test coverage,\nCI/CD pipeline", ACCENT_BLUE),
    ("Project 3", "Production Application", "pair",
     "Advanced architecture, monitoring,\nsecurity audit, evals", ACCENT_PURPLE),
]

for i, (num, title, team, desc, color) in enumerate(projects):
    x = Inches(0.8 + i * 4.0)
    # Card
    card = add_shape(slide, x, Inches(1.8), Inches(3.5), Inches(3.5), LIGHT_GRAY)
    # Color top bar
    add_shape(slide, x, Inches(1.8), Inches(3.5), Inches(0.08), color)
    # Project number
    add_textbox(slide, x + Inches(0.3), Inches(2.1), Inches(3), Inches(0.5),
                num, font_size=16, color=color, bold=True)
    # Title
    add_textbox(slide, x + Inches(0.3), Inches(2.5), Inches(3), Inches(0.6),
                title, font_size=24, color=NEU_DARK, bold=True)
    # Team type
    add_textbox(slide, x + Inches(0.3), Inches(3.1), Inches(3), Inches(0.4),
                f"({team})", font_size=16, color=MID_GRAY, bold=False)
    # Description
    add_textbox(slide, x + Inches(0.3), Inches(3.6), Inches(3), Inches(1.5),
                desc, font_size=17, color=MID_GRAY)

# Projects image alongside cards
img_path = IMG_DIR / "image_4.png"
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(0.8), Inches(5.5), Inches(5), Inches(1.8))

# Bottom callout
callout_bg = add_shape(slide, Inches(6), Inches(5.8), Inches(6.5), Inches(0.8), NEU_RED)
add_textbox(slide, Inches(6.2), Inches(5.85), Inches(6.1), Inches(0.7),
            "→  3 deployed apps you can show in interviews",
            font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: Evaluation & Format
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_red_accent_bar(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "Evaluation & Format", font_size=44, color=NEU_DARK, bold=True)

# Grading bars (horizontal stacked bar style)
components = [
    ("Projects", "50%", "P1: 13%  ·  P2: 18%  ·  P3: 19%", NEU_RED, 5.0),
    ("Homeworks", "25%", "5 scaffolding assignments", ACCENT_BLUE, 2.5),
    ("Participation", "15%", "Pre-class questions + in-class", ACCENT_GREEN, 1.5),
    ("Quizzes", "10%", "Weekly, drop lowest 2", ACCENT_PURPLE, 1.0),
]

for i, (name, pct, detail, color, bar_width) in enumerate(components):
    y = Inches(1.8 + i * 1.05)
    # Label
    add_textbox(slide, Inches(0.8), y, Inches(2.2), Inches(0.5),
                name, font_size=22, color=NEU_DARK, bold=True)
    # Bar
    add_shape(slide, Inches(3.2), y + Inches(0.05), Inches(bar_width), Inches(0.4), color)
    # Percentage on bar
    add_textbox(slide, Inches(3.3), y, Inches(1), Inches(0.5),
                pct, font_size=20, color=WHITE, bold=True)
    # Detail
    add_textbox(slide, Inches(3.2 + bar_width + 0.2), y, Inches(6), Inches(0.5),
                detail, font_size=18, color=MID_GRAY)

# Format section
add_shape(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.04), RGBColor(0xEE, 0xEE, 0xEE))

add_multiline_textbox(slide, Inches(0.8), Inches(5.9), Inches(11), Inches(1.2),
                      [
                          ("Format", True, NEU_DARK, 24),
                          "📅  Once per week  ·  3+ hours per session  ·  📍 Silicon Valley campus",
                      ],
                      font_size=22, color=MID_GRAY, line_spacing=1.5)

# ============================================================
# SLIDE 7: Prerequisites & Tools
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_red_accent_bar(slide)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "What You Need", font_size=44, color=NEU_DARK, bold=True)

# Prerequisites section
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5), Inches(0.5),
            "Prerequisites", font_size=28, color=NEU_RED, bold=True)

add_multiline_textbox(slide, Inches(0.8), Inches(2.5), Inches(5), Inches(2.5),
                      [
                          "▸  CS 5010 (min D)  or  CS 5004 (min C)",
                          "",
                          "▸  Programming fundamentals required",
                          "",
                          ("▸  No prior AI tool experience needed", True, ACCENT_GREEN, 22),
                      ],
                      font_size=22, color=MID_GRAY, line_spacing=1.3)

# Tools section
add_textbox(slide, Inches(7), Inches(1.8), Inches(5), Inches(0.5),
            "Tools (~$40/month)", font_size=28, color=NEU_RED, bold=True)

add_multiline_textbox(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(2.5),
                      [
                          "▸  Claude Pro — $20/mo",
                          "",
                          "▸  AI-Powered IDE — ~$20/mo",
                          "",
                          "▸  GitHub — free (Pro for students)",
                      ],
                      font_size=22, color=MID_GRAY, line_spacing=1.3)

# Divider between sections
add_shape(slide, Inches(6.3), Inches(1.8), Inches(0.04), Inches(3.5), RGBColor(0xEE, 0xEE, 0xEE))

# ============================================================
# SLIDE 8: Outcomes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NEU_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), NEU_RED)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "What You'll Walk Away With", font_size=44, color=WHITE, bold=True)

outcomes = [
    "3 deployed, portfolio-ready applications",
    "Professional AI-assisted development workflow",
    "Full-stack skills: React, Node.js, databases, CI/CD",
    "TDD, evals, and production engineering practices",
    "Silicon Valley–relevant tech stack proficiency",
    "Confidence to use AI tools in interviews and on the job",
]

for i, outcome in enumerate(outcomes):
    y = Inches(1.8 + i * 0.8)
    # Checkmark
    add_textbox(slide, Inches(1.0), y, Inches(0.5), Inches(0.5),
                "✓", font_size=26, color=ACCENT_GREEN, bold=True)
    # Text
    add_textbox(slide, Inches(1.6), y, Inches(10), Inches(0.5),
                outcome, font_size=24, color=RGBColor(0xDD, 0xDD, 0xDD))

# Students outcomes image on right side
img_path = IMG_DIR / "image_5.png"
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(7.5), Inches(1.5), Inches(5.5), Inches(3.1))

# ============================================================
# SLIDE 9: Call to Action
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NEU_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), NEU_RED)

# Big title
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
            "CS 7180 — Fall 2026", font_size=52, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

# CTA
add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.8),
            "Register During Course Selection!", font_size=32, color=NEU_RED, bold=True,
            alignment=PP_ALIGN.CENTER)

# Divider
add_shape(slide, Inches(5), Inches(3.3), Inches(3), Inches(0.04), NEU_RED)

# Contact info
add_multiline_textbox(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(3),
                      [
                          ("🌐  johnguerra.co/classes/aiCoding_spring_2026", False, RGBColor(0xAA, 0xCC, 0xFF), 24),
                          "",
                          ("📧  jguerra@northeastern.edu", False, RGBColor(0xCC, 0xCC, 0xCC), 24),
                          "",
                          ("John Alexis Guerra Gomez", True, WHITE, 24),
                          ("Khoury College of Computer Sciences", False, RGBColor(0xAA, 0xAA, 0xAA), 20),
                          ("Northeastern University — Silicon Valley", False, RGBColor(0xAA, 0xAA, 0xAA), 20),
                      ],
                      font_size=24, color=RGBColor(0xCC, 0xCC, 0xCC), line_spacing=1.3,
                      alignment=PP_ALIGN.CENTER)

# Husky/circuit background image
img_path = IMG_DIR / "image_6.png"
if img_path.exists():
    pic = slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    # Send to back
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

# Disclaimer on last slide
add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
            "Slides generated with Claude Code · Images generated with Gemini",
            font_size=12, color=RGBColor(0x66, 0x66, 0x77), alignment=PP_ALIGN.CENTER)

# Save
output_path = "/Users/aguerra/workspace/aiCoding_Course/other/preview-night/CS7180_Preview_Night_Fall2026.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
